"""300kW 수소엔진 최종 발표자료 PPTX v2 - 5대 핵심기술 강조"""
import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 색상
DARK_BLUE = RGBColor(0x00, 0x2B, 0x5C)
MEDIUM_BLUE = RGBColor(0x00, 0x4E, 0x92)
LIGHT_BLUE = RGBColor(0x00, 0x7B, 0xC0)
ACCENT_ORANGE = RGBColor(0xE8, 0x6C, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MID_GRAY = RGBColor(0x66, 0x66, 0x66)
TABLE_HEADER_BG = RGBColor(0x00, 0x2B, 0x5C)
TABLE_ALT_BG = RGBColor(0xE8, 0xF0, 0xF8)
TABLE_WHITE_BG = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x00, 0x7A, 0x33)
RED_ACCENT = RGBColor(0xC0, 0x39, 0x2B)
HIGHLIGHT_YELLOW = RGBColor(0xFF, 0xF3, 0xCD)
HIGHLIGHT_GREEN = RGBColor(0xD4, 0xED, 0xDA)
HIGHLIGHT_BLUE = RGBColor(0xCC, 0xE5, 0xFF)


def add_title_bar(slide, title_text, subtitle_text=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                   Inches(13.333), Inches(1.1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12), Inches(0.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if subtitle_text:
        txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.7), Inches(12), Inches(0.35))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle_text
        p2.font.size = Pt(14)
        p2.font.color.rgb = RGBColor(0xAA, 0xCC, 0xEE)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.1),
                                  Inches(13.333), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_ORANGE
    line.line.fill.background()


def add_footer(slide, page_num, total=3):
    txBox = slide.shapes.add_textbox(Inches(11.5), Inches(7.1), Inches(1.5), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{page_num} / {total}"
    p.font.size = Pt(10)
    p.font.color.rgb = MID_GRAY
    p.alignment = PP_ALIGN.RIGHT
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(4), Inches(0.3))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "한국기계연구원(KIMM) | 산업기술혁신사업"
    p2.font.size = Pt(10)
    p2.font.color.rgb = MID_GRAY


def set_cell(cell, text, sz=10, bold=False, color=DARK_GRAY, align=PP_ALIGN.LEFT, bg=None):
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if bg:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg


def add_box(slide, left, top, width, height, text, bg, sz=11, tc=WHITE, bold=True,
            align=PP_ALIGN.CENTER, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(6)
    tf.margin_top = Pt(3)
    tf.margin_bottom = Pt(3)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.bold = bold
    p.font.color.rgb = tc
    p.alignment = align
    return shape


def add_text(slide, left, top, width, height, text, sz=11, bold=False, color=DARK_GRAY,
             align=PP_ALIGN.LEFT, line_sp=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    if line_sp:
        p.line_spacing = Pt(line_sp)
    return txBox


# ============================================================
# 슬라이드 1: 연차별 요약 (기존과 동일)
# ============================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
slide1.background.fill.solid()
slide1.background.fill.fore_color.rgb = WHITE
add_title_bar(slide1, "300kW급 수소 엔진 개발 - 연차별 추진 요약",
              "공동연구개발기관: 한국기계연구원 | 1단계 1~4차년도")
add_footer(slide1, 1)

card_data = [
    {"year": "1차년도", "theme": "시스템 구축 · 기초 성능",
     "items": ["동력계 엔진 성능 평가\n시스템 구축 (460kW AC동력계)",
               "수소 연료 공급 설비 증축\n(35kg 저장, 60bar 대응)",
               "베이스 천연가스 엔진\n(GX12) 성능 DB 구축",
               "수소 엔진 기초 시험\n→ 열효율 2.8%p 향상"],
     "highlight": "열효율 2.8%p 향상",
     "color": RGBColor(0x00, 0x6B, 0x9F)},
    {"year": "2차년도", "theme": "하드웨어 선정 · 요소기술",
     "items": ["냉형 점화플러그 적용\n→ 이른 조기점화 억제",
               "포트분사 인젝터 선정\n(분사 인가기간 <10ms)",
               "최적 압축비 11:1 선정\n→ 열효율 ~1%p 향상",
               "피스톤 형상 Type C 제안\n→ 이상연소 동시 억제"],
     "highlight": "압축비 11:1 + Type C",
     "color": RGBColor(0x00, 0x4E, 0x92)},
    {"year": "3차년도", "theme": "NOx 저감 · 후처리시스템",
     "items": ["NOx 저감 연소 제어 전략\n(분사/점화시기, 공기과잉률)",
               "저부하 EGR 전략\n→ 펌핑손실 저감",
               "후처리시스템 간소화\nDOC + SCR (DPF 삭제)",
               "엔진-후처리 연계 시험\n→ Euro 7 배출규제 충족"],
     "highlight": "Euro 7 충족 달성",
     "color": RGBColor(0x2E, 0x86, 0xAB)},
    {"year": "4차년도", "theme": "통합제어 · 고성능화",
     "items": ["운전영역별 연소-후처리\n통합제어 전략 도출",
               "주행 사이클 대응\n(WHSC/WHTC/NRSC/NRTC)",
               "2-stage 터보차저 적용\n→ 저속토크 42.6% 향상",
               "고성능 수소엔진\n개발 로드맵 제시"],
     "highlight": "저속토크 42.6% 향상",
     "color": RGBColor(0x00, 0x2B, 0x5C)},
]

cw = Inches(3.0)
cg = Inches(0.15)
sx = Inches(0.45)
ct = Inches(1.4)

for i, c in enumerate(card_data):
    x = sx + i * (cw + cg)
    add_box(slide1, x, ct, cw, Inches(0.45), c["year"], c["color"], sz=14)
    add_box(slide1, x, ct + Inches(0.5), cw, Inches(0.35), c["theme"],
            RGBColor(0xEE, 0xEE, 0xEE), sz=11, tc=DARK_GRAY)
    for j, item in enumerate(c["items"]):
        box = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      x + Inches(0.05), ct + Inches(1.0) + j * Inches(1.15),
                                      cw - Inches(0.1), Inches(1.05))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_GRAY
        box.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
        box.line.width = Pt(0.5)
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(8)
        tf.margin_right = Pt(8)
        tf.margin_top = Pt(5)
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(10)
        p.font.color.rgb = DARK_GRAY
        p.line_spacing = Pt(14)
    ht = ct + Inches(1.0) + 4 * Inches(1.15) + Inches(0.1)
    add_box(slide1, x, ht, cw, Inches(0.4), f"★ {c['highlight']}", ACCENT_ORANGE, sz=11)


# ============================================================
# 슬라이드 2: 핵심결과 (1) - 핵심 요소기술 3가지
# ============================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
slide2.background.fill.solid()
slide2.background.fill.fore_color.rgb = WHITE
add_title_bar(slide2, "핵심결과 (1) - 핵심 요소기술 개발",
              "냉형 점화플러그 · 포트분사 인젝터 · 피스톤 형상 최적화")
add_footer(slide2, 2)

# --- 좌측: 냉형 점화플러그 ---
add_text(slide2, Inches(0.5), Inches(1.3), Inches(4), Inches(0.35),
         "1  냉형 점화플러그 적용", sz=15, bold=True, color=DARK_BLUE)

sp_table = slide2.shapes.add_table(4, 3, Inches(0.5), Inches(1.75),
                                   Inches(5.8), Inches(1.85)).table
sp_table.columns[0].width = Inches(1.5)
sp_table.columns[1].width = Inches(2.15)
sp_table.columns[2].width = Inches(2.15)

for j, h in enumerate(["항목", "베이스 점화플러그", "냉형 점화플러그"]):
    set_cell(sp_table.cell(0, j), h, sz=10, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, bg=TABLE_HEADER_BG)

sp_data = [
    ["이른 조기점화", "고부하 빈번 발생\n→ 순간 출력 하락", "완전 억제"],
    ["늦은 조기점화", "다수 발생", "발생빈도 대폭 감소"],
    ["100% 부하 운전", "운전 제한", "고부하 안정 운전 확보"],
]
for i, row in enumerate(sp_data):
    bg = TABLE_ALT_BG if i % 2 == 0 else TABLE_WHITE_BG
    set_cell(sp_table.cell(i+1, 0), row[0], sz=9, bold=True, bg=bg)
    set_cell(sp_table.cell(i+1, 1), row[1], sz=9, bg=bg, color=RED_ACCENT)
    set_cell(sp_table.cell(i+1, 2), row[2], sz=9, bg=bg, color=GREEN, bold=True)

add_box(slide2, Inches(0.5), Inches(3.7), Inches(5.8), Inches(0.4),
        "점화플러그 온도 저감 → λ 2.0~2.2에서 이른 조기점화 완전 억제",
        HIGHLIGHT_GREEN, sz=10, tc=GREEN, line_color=GREEN)

# --- 우측: 포트분사 인젝터 ---
add_text(slide2, Inches(7.0), Inches(1.3), Inches(5.5), Inches(0.35),
         "2  포트분사 인젝터 사양 선정", sz=15, bold=True, color=DARK_BLUE)

inj_table = slide2.shapes.add_table(4, 3, Inches(7.0), Inches(1.75),
                                    Inches(5.8), Inches(1.85)).table
inj_table.columns[0].width = Inches(1.5)
inj_table.columns[1].width = Inches(2.15)
inj_table.columns[2].width = Inches(2.15)

for j, h in enumerate(["평가 항목", "작은 구경", "큰 구경 (선정)"]):
    set_cell(inj_table.cell(0, j), h, sz=10, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, bg=TABLE_HEADER_BG)

inj_data = [
    ["아이들 안정성", "유리", "COV of IMEP < 3%"],
    ["분사 인가기간", "과도 → 이상연소", "< 10ms 만족"],
    ["300kW 대응", "불가", "충분한 연료 공급"],
]
for i, row in enumerate(inj_data):
    bg = TABLE_ALT_BG if i % 2 == 0 else TABLE_WHITE_BG
    set_cell(inj_table.cell(i+1, 0), row[0], sz=9, bold=True, bg=bg)
    set_cell(inj_table.cell(i+1, 1), row[1], sz=9, bg=bg, color=RED_ACCENT)
    set_cell(inj_table.cell(i+1, 2), row[2], sz=9, bg=bg, color=GREEN, bold=True)

add_box(slide2, Inches(7.0), Inches(3.7), Inches(5.8), Inches(0.4),
        "정격출력 300kW + 아이들 안정성 동시 만족하는 큰 구경 인젝터 선정",
        HIGHLIGHT_GREEN, sz=10, tc=GREEN, line_color=GREEN)

# --- 하단: 피스톤 형상 최적화 ---
add_text(slide2, Inches(0.5), Inches(4.35), Inches(12), Inches(0.35),
         "3  피스톤 형상 최적화 (압축비 + 연소실 형상)", sz=15, bold=True, color=DARK_BLUE)

pist_table = slide2.shapes.add_table(5, 5, Inches(0.5), Inches(4.8),
                                     Inches(12.3), Inches(2.0)).table
pist_table.columns[0].width = Inches(1.8)
pist_table.columns[1].width = Inches(1.2)
pist_table.columns[2].width = Inches(1.5)
pist_table.columns[3].width = Inches(4.0)
pist_table.columns[4].width = Inches(3.8)

for j, h in enumerate(["사양", "압축비", "Squish", "특성", "열효율 성과"]):
    set_cell(pist_table.cell(0, j), h, sz=10, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, bg=TABLE_HEADER_BG)

pist_data = [
    ["Base (GX12)", "10:1", "기본", "천연가스 엔진 기준 설계", "기준선"],
    ["Type A", "11:1", "기본", "압축비 증대 → 열효율 향상", "~1%p 향상, 조기점화 없음"],
    ["Type B", "11:1", "좁음", "화염전파 유리, 고부하 노킹 취약", "연소효율 소폭 증가"],
    ["Type C (최종)", "11:1", "넓음", "Squish 유동 강화 → 조기점화+노킹 동시 억제",
     "40.4% (1800rpm) / 41.4% (1200rpm)"],
]
for i, row in enumerate(pist_data):
    is_final = (i == 3)
    bg_r = HIGHLIGHT_YELLOW if is_final else (TABLE_ALT_BG if i % 2 == 0 else TABLE_WHITE_BG)
    for j, val in enumerate(row):
        c = ACCENT_ORANGE if is_final else DARK_GRAY
        set_cell(pist_table.cell(i+1, j), val, sz=9, bold=is_final, bg=bg_r, color=c,
                 align=PP_ALIGN.CENTER if j < 3 else PP_ALIGN.LEFT)


# ============================================================
# 슬라이드 3: 핵심결과 (2) - 통합제어 + 2-stage TC
# ============================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
slide3.background.fill.solid()
slide3.background.fill.fore_color.rgb = WHITE
add_title_bar(slide3, "핵심결과 (2) - 통합제어 전략 및 고성능화",
              "운전영역별 연소-후처리 통합제어 · 2-Stage 터보차저 · 향후 로드맵")
add_footer(slide3, 3)

# --- 좌측: 통합제어 전략 ---
add_text(slide3, Inches(0.5), Inches(1.3), Inches(6), Inches(0.35),
         "4  운전영역별 연소-후처리 통합제어 전략", sz=15, bold=True, color=DARK_BLUE)

# 제어변수 영향
ctrl_table = slide3.shapes.add_table(5, 4, Inches(0.5), Inches(1.75),
                                     Inches(6.0), Inches(2.2)).table
ctrl_table.columns[0].width = Inches(1.3)
ctrl_table.columns[1].width = Inches(1.7)
ctrl_table.columns[2].width = Inches(1.5)
ctrl_table.columns[3].width = Inches(1.5)

for j, h in enumerate(["운전영역", "열효율 지배인자", "NOx 제어", "제한 요소"]):
    set_cell(ctrl_table.cell(0, j), h, sz=9, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, bg=TABLE_HEADER_BG)

ctrl_data = [
    ["중속 저부하", "공기과잉률 지배", "높은 λ", "-"],
    ["중속 중/고부하", "연소상+λ 동시", "높은 λ+지각", "NOx"],
    ["고속 저/중부하", "연소상+λ 동시", "낮은 λ 가능", "-"],
    ["고속 고부하", "연소상 지배", "낮은 λ+지각", "노킹/조기점화"],
]
for i, row in enumerate(ctrl_data):
    bg = TABLE_ALT_BG if i % 2 == 0 else TABLE_WHITE_BG
    for j, val in enumerate(row):
        set_cell(ctrl_table.cell(i+1, j), val, sz=9, bg=bg,
                 bold=(j == 0), align=PP_ALIGN.CENTER)

# 최종 전략 요약
add_text(slide3, Inches(0.5), Inches(4.1), Inches(6), Inches(0.3),
         "최종 전략 요약", sz=12, bold=True, color=DARK_BLUE)

strat_table = slide3.shapes.add_table(3, 3, Inches(0.5), Inches(4.45),
                                      Inches(6.0), Inches(1.0)).table
strat_table.columns[0].width = Inches(1.0)
strat_table.columns[1].width = Inches(2.5)
strat_table.columns[2].width = Inches(2.5)

for j, h in enumerate(["부하", "중속", "고속"]):
    set_cell(strat_table.cell(0, j), h, sz=10, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, bg=TABLE_HEADER_BG)

strat_data = [
    ["저부하", "높은 λ + 최적 연소상\n→ 최고 열효율", "낮은 λ + 최적 연소상"],
    ["고부하", "높은 λ + 지각 연소상\n→ NOx 억제", "낮은 λ + 지각 연소상\n→ 출력 확보"],
]
for i, row in enumerate(strat_data):
    bg = HIGHLIGHT_BLUE if i == 0 else HIGHLIGHT_YELLOW
    for j, val in enumerate(row):
        set_cell(strat_table.cell(i+1, j), val, sz=9, bg=bg,
                 bold=(j == 0), align=PP_ALIGN.CENTER)

# 후처리 간소화
add_box(slide3, Inches(0.5), Inches(5.6), Inches(2.9), Inches(0.55),
        "디젤: DOC→DPF→SCR\n(3단 구성)", MID_GRAY, sz=10)

add_text(slide3, Inches(3.45), Inches(5.65), Inches(0.4), Inches(0.4),
         "→", sz=20, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

add_box(slide3, Inches(3.75), Inches(5.6), Inches(2.75), Inches(0.55),
        "수소: DOC→SCR\n(2단, DPF삭제+사이즈축소)", GREEN, sz=10)

# Euro 7
add_box(slide3, Inches(0.5), Inches(6.25), Inches(6.0), Inches(0.45),
        "TC후단 400℃↑ → Urea SCR 99% 변환효율 → Euro 7 충족",
        HIGHLIGHT_GREEN, sz=11, tc=GREEN, bold=True, line_color=GREEN)

# --- 우측: 2-stage TC ---
add_text(slide3, Inches(7.0), Inches(1.3), Inches(5.5), Inches(0.35),
         "5  2-Stage 터보차저 적용 성과", sz=15, bold=True, color=DARK_BLUE)

tc_table = slide3.shapes.add_table(4, 4, Inches(7.0), Inches(1.75),
                                   Inches(5.8), Inches(1.6)).table
tc_table.columns[0].width = Inches(0.9)
tc_table.columns[1].width = Inches(1.2)
tc_table.columns[2].width = Inches(1.7)
tc_table.columns[3].width = Inches(2.0)

for j, h in enumerate(["속도", "토크 향상", "제한 요소", "최적화 전략"]):
    set_cell(tc_table.cell(0, j), h, sz=10, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, bg=TABLE_HEADER_BG)

tc_data = [
    ["저속", "+42.6%", "조기점화\n배기에너지 부족", "점화시기 지각\n+ 적정 당량비"],
    ["중속", "+6.1%", "인터쿨러 전단온도\nNOx 배출", "고정 과급압에서\n당량비 최적화"],
    ["고속", "+3.4%", "인터쿨러 전단온도\n조기점화", "NOx 마진 활용\n당량비 상향"],
]
for i, row in enumerate(tc_data):
    is_first = (i == 0)
    bg = HIGHLIGHT_YELLOW if is_first else (TABLE_ALT_BG if i % 2 == 0 else TABLE_WHITE_BG)
    set_cell(tc_table.cell(i+1, 0), row[0], sz=10, bold=True, bg=bg, align=PP_ALIGN.CENTER)
    set_cell(tc_table.cell(i+1, 1), row[1], sz=14 if is_first else 11, bold=True,
             bg=bg, color=ACCENT_ORANGE, align=PP_ALIGN.CENTER)
    set_cell(tc_table.cell(i+1, 2), row[2], sz=9, bg=bg, align=PP_ALIGN.CENTER)
    set_cell(tc_table.cell(i+1, 3), row[3], sz=9, bg=bg, align=PP_ALIGN.CENTER)

# 저속 토크 메커니즘
add_text(slide3, Inches(7.0), Inches(3.5), Inches(5.5), Inches(0.3),
         "저속 토크 향상 메커니즘", sz=12, bold=True, color=DARK_BLUE)

mech_box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(7.0), Inches(3.85), Inches(5.8), Inches(1.5))
mech_box.fill.solid()
mech_box.fill.fore_color.rgb = LIGHT_GRAY
mech_box.line.color.rgb = MEDIUM_BLUE
mech_box.line.width = Pt(1)
tf = mech_box.text_frame
tf.word_wrap = True
tf.margin_left = Pt(10)
tf.margin_right = Pt(10)
tf.margin_top = Pt(8)

lines = [
    "▸ 점화시기 지각 → 배기에너지 확보 → 과급압 형성 (가장 효과적)",
    "   └ 한계: 실린더 압력/온도 상승 → 조기점화 발생",
    "▸ 당량비 최적화:",
    "   └ 낮은 당량비: VGT duty↑ → 배기압↑ → 잔여가스↑ → 조기점화",
    "   └ 높은 당량비: 연소온도↑ → 조기점화 재발생",
    "   └ 적정 당량비에서 조기점화 최소 + 열효율 최적",
]
p = tf.paragraphs[0]
p.text = lines[0]
p.font.size = Pt(9)
p.font.color.rgb = DARK_GRAY
p.line_spacing = Pt(13)
for line in lines[1:]:
    p2 = tf.add_paragraph()
    p2.text = line
    p2.font.size = Pt(9)
    p2.font.color.rgb = DARK_GRAY
    p2.line_spacing = Pt(13)

# 향후 로드맵
add_text(slide3, Inches(7.0), Inches(5.5), Inches(5.5), Inches(0.3),
         "향후 고성능화 로드맵", sz=12, bold=True, color=DARK_BLUE)

roadmap = [
    ("하이브리드화", "저속토크/반응성 보완", MEDIUM_BLUE),
    ("2단 인터쿨러", "중/고속 토크 추가 향상", RGBColor(0x2E, 0x86, 0xAB)),
    ("Blowby 관리", "5%↑ H₂ 열점 관리", RGBColor(0x8E, 0x44, 0xAD)),
]
for i, (title, desc, color) in enumerate(roadmap):
    x_pos = Inches(7.0) + i * Inches(1.95)
    add_box(slide3, x_pos, Inches(5.85), Inches(1.85), Inches(0.35),
            title, color, sz=10)
    add_box(slide3, x_pos, Inches(6.25), Inches(1.85), Inches(0.45),
            desc, LIGHT_GRAY, sz=9, tc=DARK_GRAY, bold=False, line_color=color)


# 저장
output_path = r"C:\Users\hwpar\OneDrive\퇴직준비세미나\Projects\요약자료\output\300kW_수소엔진_발표자료_v2.pptx"
prs.save(output_path)
print(f"PPTX v2 생성 완료: {output_path}")
