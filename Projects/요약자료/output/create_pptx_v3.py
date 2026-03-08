"""300kW 수소엔진 발표자료 PPTX v3 - 상세 연차요약 + NOx 전략 추가"""
import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

DARK_BLUE = RGBColor(0x00, 0x2B, 0x5C)
MEDIUM_BLUE = RGBColor(0x00, 0x4E, 0x92)
ACCENT_ORANGE = RGBColor(0xE8, 0x6C, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MID_GRAY = RGBColor(0x66, 0x66, 0x66)
TH_BG = RGBColor(0x00, 0x2B, 0x5C)
ALT_BG = RGBColor(0xE8, 0xF0, 0xF8)
W_BG = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x00, 0x7A, 0x33)
RED = RGBColor(0xC0, 0x39, 0x2B)
HL_Y = RGBColor(0xFF, 0xF3, 0xCD)
HL_G = RGBColor(0xD4, 0xED, 0xDA)
HL_B = RGBColor(0xCC, 0xE5, 0xFF)
PURPLE = RGBColor(0x6C, 0x3C, 0x97)


def title_bar(s, t, sub=None):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.1))
    sh.fill.solid(); sh.fill.fore_color.rgb = DARK_BLUE; sh.line.fill.background()
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12), Inches(0.6))
    p = tb.text_frame.paragraphs[0]; p.text = t; p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = WHITE
    if sub:
        tb2 = s.shapes.add_textbox(Inches(0.5), Inches(0.7), Inches(12), Inches(0.35))
        p2 = tb2.text_frame.paragraphs[0]; p2.text = sub; p2.font.size = Pt(14); p2.font.color.rgb = RGBColor(0xAA, 0xCC, 0xEE)
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.1), Inches(13.333), Inches(0.04))
    ln.fill.solid(); ln.fill.fore_color.rgb = ACCENT_ORANGE; ln.line.fill.background()


def footer(s, n, total=3):
    tb = s.shapes.add_textbox(Inches(11.5), Inches(7.1), Inches(1.5), Inches(0.3))
    p = tb.text_frame.paragraphs[0]; p.text = f"{n} / {total}"; p.font.size = Pt(10); p.font.color.rgb = MID_GRAY; p.alignment = PP_ALIGN.RIGHT
    tb2 = s.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(4), Inches(0.3))
    p2 = tb2.text_frame.paragraphs[0]; p2.text = "한국기계연구원(KIMM) | 산업기술혁신사업"; p2.font.size = Pt(10); p2.font.color.rgb = MID_GRAY


def cell(c, t, sz=10, b=False, co=DARK_GRAY, a=PP_ALIGN.LEFT, bg=None):
    c.text = ""; p = c.text_frame.paragraphs[0]; p.text = t; p.font.size = Pt(sz); p.font.bold = b; p.font.color.rgb = co; p.alignment = a
    c.vertical_anchor = MSO_ANCHOR.MIDDLE
    if bg: c.fill.solid(); c.fill.fore_color.rgb = bg


def box(s, l, t, w, h, txt, bg, sz=11, tc=WHITE, b=True, a=PP_ALIGN.CENTER, lc=None):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = bg
    if lc: sh.line.color.rgb = lc; sh.line.width = Pt(1.5)
    else: sh.line.fill.background()
    tf = sh.text_frame; tf.word_wrap = True; tf.margin_left = Pt(6); tf.margin_right = Pt(6); tf.margin_top = Pt(3); tf.margin_bottom = Pt(3)
    p = tf.paragraphs[0]; p.text = txt; p.font.size = Pt(sz); p.font.bold = b; p.font.color.rgb = tc; p.alignment = a
    return sh


def txt(s, l, t, w, h, text, sz=11, b=False, co=DARK_GRAY, a=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(l, t, w, h); tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(sz); p.font.bold = b; p.font.color.rgb = co; p.alignment = a
    return tb


def multiline_box(s, l, t, w, h, lines, bg, lc, sz=9, co=DARK_GRAY):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = bg
    sh.line.color.rgb = lc; sh.line.width = Pt(1)
    tf = sh.text_frame; tf.word_wrap = True; tf.margin_left = Pt(10); tf.margin_right = Pt(10); tf.margin_top = Pt(6)
    p = tf.paragraphs[0]; p.text = lines[0]; p.font.size = Pt(sz); p.font.color.rgb = co; p.line_spacing = Pt(13)
    for line in lines[1:]:
        p2 = tf.add_paragraph(); p2.text = line; p2.font.size = Pt(sz); p2.font.color.rgb = co; p2.line_spacing = Pt(13)
    return sh


# ============================================================
# 슬라이드 1: 연차별 추진 요약 (상세 3~4줄)
# ============================================================
s1 = prs.slides.add_slide(prs.slide_layouts[6])
s1.background.fill.solid(); s1.background.fill.fore_color.rgb = WHITE
title_bar(s1, "300kW급 수소 엔진 개발 - 연차별 추진 요약", "공동연구개발기관: 한국기계연구원 | 1단계 1~4차년도")
footer(s1, 1)

years = [
    {"year": "1차년도", "theme": "시스템 구축 · 기초 성능 확보",
     "color": RGBColor(0x00, 0x6B, 0x9F),
     "lines": [
         "460kW AC동력계 기반 엔진 성능 평가 시스템 구축, 배출가스 분석기(AMA i60), PM 분석기, 코리올리식 질량유량계 등 계측장비 확보",
         "수소 저장고 증축(랙 타입 70여개 용기, 35kg), 엔진셀까지 고압배관(60bar 대응), 연료 공급 및 안전 설비 구축 완료",
         "베이스 천연가스 엔진(GX12) 운전속도·부하별 성능/배출가스/연소인자 DB 구축",
         "수소 엔진 기초 시험에서 천연가스 대비 제동열효율 2.8%p 향상 달성 (펌핑손실·열전달 손실 저감 효과)",
     ]},
    {"year": "2차년도", "theme": "하드웨어 선정 · 요소기술 평가",
     "color": RGBColor(0x00, 0x4E, 0x92),
     "lines": [
         "냉형 점화플러그 적용으로 100% 부하에서 이른 조기점화 완전 억제, 백금 제거 사양에서도 유사한 효과 확인",
         "정격 300kW 달성에 필요한 분사 인가기간(<10ms)과 아이들 안정성(COV <3%)을 동시 만족하는 큰 구경 포트분사 인젝터 선정",
         "단기통 엔진에서 공기과잉률 및 EGR 전략 평가 → 이상연소 억제 효과적이나 터보차저 부담과의 trade-off 규명",
         "최적 압축비 11:1 선정(열효율 ~1%p↑), Squish 유동 강화 피스톤 Type C 제안(조기점화+노킹 동시 억제)",
     ]},
    {"year": "3차년도", "theme": "NOx 저감 · 후처리시스템 구성",
     "color": RGBColor(0x2E, 0x86, 0xAB),
     "lines": [
         "분사시기 지연 시 국부 농후 혼합기에 의한 NOx 증가 확인 → 이른 분사시기+높은 분사압력+지각 점화시기 조합 도출",
         "저부하 EGR 적용 시 펌핑손실 저감 → 열효율 향상, 과잉 공기 대비 낮은 λ에서 NOx 저감 가능 (Transient 유리)",
         "CO/CO₂/PM 제로 활용하여 DPF 삭제 → DOC+SCR 2단 구성으로 후처리 간소화, 사이즈 축소 가능성 제시",
         "엔진-후처리 연계 시험: 25~100% 부하 전영역 TC후단 400℃↑, SCR 변환효율 99% → Euro 7 충족",
     ]},
    {"year": "4차년도", "theme": "통합제어 · 고성능화",
     "color": RGBColor(0x00, 0x2B, 0x5C),
     "lines": [
         "운전영역별(중속/고속 × 저/중/고부하) 열효율 지배인자 분석, NOx 마진·이상연소 한계 고려 연소-후처리 통합제어 전략 도출",
         "WHSC/WHTC/NRSC/NRTC 주행사이클 시험 탑재, Transient 운전 고려 제어변수 최적화(조기점화·노킹 대응)",
         "2-stage 터보차저 적용: 저속토크 +42.6%, 중속 +6.1%, 고속 +3.4% 향상 (점화시기 지각+적정 당량비 선정 핵심)",
         "하이브리드화(저속토크 보완), 2단 인터쿨러(중/고속 토크 향상), Blowby 열점 관리 등 향후 로드맵 제시",
     ]},
]

for i, y in enumerate(years):
    col_x = Inches(0.4) + i * Inches(3.2)
    col_w = Inches(3.05)
    # 연차 헤더
    box(s1, col_x, Inches(1.35), col_w, Inches(0.42), y["year"], y["color"], sz=14)
    # 테마
    box(s1, col_x, Inches(1.82), col_w, Inches(0.32), y["theme"],
        RGBColor(0xEE, 0xEE, 0xEE), sz=10, tc=DARK_GRAY)

    # 상세 내용 (4줄)
    for j, line in enumerate(y["lines"]):
        item_y = Inches(2.25) + j * Inches(1.25)
        bx = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 col_x + Inches(0.03), item_y,
                                 col_w - Inches(0.06), Inches(1.15))
        bx.fill.solid()
        bx.fill.fore_color.rgb = LIGHT_GRAY
        bx.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
        bx.line.width = Pt(0.5)
        tf = bx.text_frame; tf.word_wrap = True
        tf.margin_left = Pt(7); tf.margin_right = Pt(7); tf.margin_top = Pt(4)
        p = tf.paragraphs[0]
        p.text = line
        p.font.size = Pt(8)
        p.font.color.rgb = DARK_GRAY
        p.line_spacing = Pt(11)


# ============================================================
# 슬라이드 2: 핵심결과 (1) - 요소기술 + NOx 전략
# ============================================================
s2 = prs.slides.add_slide(prs.slide_layouts[6])
s2.background.fill.solid(); s2.background.fill.fore_color.rgb = WHITE
title_bar(s2, "핵심결과 (1) - 핵심 요소기술 및 NOx 저감 전략",
          "냉형 점화플러그 · 포트분사 인젝터 · 피스톤 형상 · NOx 연소제어")
footer(s2, 2)

# --- 좌상: 냉형 점화플러그 ---
txt(s2, Inches(0.4), Inches(1.25), Inches(4), Inches(0.3), "1  냉형 점화플러그", sz=14, b=True, co=DARK_BLUE)
t1 = s2.shapes.add_table(4, 3, Inches(0.4), Inches(1.6), Inches(6.1), Inches(1.5)).table
t1.columns[0].width = Inches(1.4); t1.columns[1].width = Inches(2.35); t1.columns[2].width = Inches(2.35)
for j, h in enumerate(["항목", "베이스", "냉형 (선정)"]):
    cell(t1.cell(0, j), h, sz=9, b=True, co=WHITE, a=PP_ALIGN.CENTER, bg=TH_BG)
for i, (a, b, c) in enumerate([
    ("이른 조기점화", "고부하 빈번→출력하락", "완전 억제"),
    ("늦은 조기점화", "다수 발생", "대폭 감소"),
    ("100% 부하", "운전 제한", "안정 운전 확보"),
]):
    bg = ALT_BG if i % 2 == 0 else W_BG
    cell(t1.cell(i+1, 0), a, sz=8, b=True, bg=bg)
    cell(t1.cell(i+1, 1), b, sz=8, bg=bg, co=RED)
    cell(t1.cell(i+1, 2), c, sz=8, bg=bg, co=GREEN, b=True)

# --- 우상: 포트분사 인젝터 ---
txt(s2, Inches(6.9), Inches(1.25), Inches(5), Inches(0.3), "2  포트분사 인젝터 선정", sz=14, b=True, co=DARK_BLUE)
t2 = s2.shapes.add_table(4, 3, Inches(6.9), Inches(1.6), Inches(6.0), Inches(1.5)).table
t2.columns[0].width = Inches(1.4); t2.columns[1].width = Inches(2.3); t2.columns[2].width = Inches(2.3)
for j, h in enumerate(["항목", "작은 구경", "큰 구경 (선정)"]):
    cell(t2.cell(0, j), h, sz=9, b=True, co=WHITE, a=PP_ALIGN.CENTER, bg=TH_BG)
for i, (a, b, c) in enumerate([
    ("아이들 안정성", "유리", "COV of IMEP <3%"),
    ("분사 인가기간", "과도→이상연소", "<10ms 만족"),
    ("300kW 대응", "불가", "충분한 연료공급"),
]):
    bg = ALT_BG if i % 2 == 0 else W_BG
    cell(t2.cell(i+1, 0), a, sz=8, b=True, bg=bg)
    cell(t2.cell(i+1, 1), b, sz=8, bg=bg, co=RED)
    cell(t2.cell(i+1, 2), c, sz=8, bg=bg, co=GREEN, b=True)

# --- 좌하: 피스톤 형상 ---
txt(s2, Inches(0.4), Inches(3.3), Inches(6), Inches(0.3), "3  피스톤 형상 최적화", sz=14, b=True, co=DARK_BLUE)
t3 = s2.shapes.add_table(5, 5, Inches(0.4), Inches(3.65), Inches(6.1), Inches(1.8)).table
t3.columns[0].width = Inches(1.3); t3.columns[1].width = Inches(0.7); t3.columns[2].width = Inches(0.7)
t3.columns[3].width = Inches(1.8); t3.columns[4].width = Inches(1.6)
for j, h in enumerate(["사양", "CR", "Squish", "특성", "열효율"]):
    cell(t3.cell(0, j), h, sz=8, b=True, co=WHITE, a=PP_ALIGN.CENTER, bg=TH_BG)
pist = [
    ("Base", "10:1", "기본", "천연가스 기준", "기준선"),
    ("Type A", "11:1", "기본", "압축비 증대", "~1%p↑"),
    ("Type B", "11:1", "좁음", "화염전파↑, 고부하 노킹", "소폭 증가"),
    ("Type C", "11:1", "넓음", "Squish 강화→이상연소 억제", "40.4%/41.4%"),
]
for i, row in enumerate(pist):
    is_c = (i == 3)
    bg = HL_Y if is_c else (ALT_BG if i % 2 == 0 else W_BG)
    co = ACCENT_ORANGE if is_c else DARK_GRAY
    for j, v in enumerate(row):
        cell(t3.cell(i+1, j), v, sz=8, b=is_c, bg=bg, co=co, a=PP_ALIGN.CENTER)

# --- 우하: NOx 저감 전략 (신규) ---
txt(s2, Inches(6.9), Inches(3.3), Inches(5.5), Inches(0.3), "4  NOx 저감 연소 제어 전략", sz=14, b=True, co=DARK_BLUE)

# 제어변수별 NOx 테이블
t4 = s2.shapes.add_table(6, 3, Inches(6.9), Inches(3.65), Inches(6.0), Inches(2.1)).table
t4.columns[0].width = Inches(1.8); t4.columns[1].width = Inches(1.6); t4.columns[2].width = Inches(2.6)
for j, h in enumerate(["제어변수", "NOx 영향", "비고"]):
    cell(t4.cell(0, j), h, sz=8, b=True, co=WHITE, a=PP_ALIGN.CENTER, bg=TH_BG)
nox_data = [
    ("공기과잉률 증가", "대폭 감소", "역화/조기점화 억제, but 부하↓"),
    ("점화시기 지각", "감소", "노킹 억제, 열효율 소폭↓"),
    ("분사시기 진각", "감소 (균일혼합기)", "조기점화 억제"),
    ("분사시기 지각", "증가 (국부농후)", "조기점화 유발 → 지양"),
    ("EGR 적용", "감소", "펌핑손실↓, 열효율↑ (저부하)"),
]
for i, (a, b, c) in enumerate(nox_data):
    bg = ALT_BG if i % 2 == 0 else W_BG
    is_bad = ("증가" in b and "감소" not in b)
    cell(t4.cell(i+1, 0), a, sz=8, b=True, bg=bg)
    cell(t4.cell(i+1, 1), b, sz=8, bg=bg, co=RED if is_bad else GREEN, b=True)
    cell(t4.cell(i+1, 2), c, sz=8, bg=bg)

# NOx 부하별 전략 요약
nox_summary = [
    "▸ 저부하(25%): 배기온도 낮아 SCR 효율↓ → 엔진측 NOx 저감 필수 (높은 λ or EGR)",
    "▸ 중부하(50~75%): 공기과잉률 + 점화시기 지각, TC후단 400℃↑ → SCR 99%",
    "▸ 고부하(100%): 출력 우선(λ 증가 제한), 점화시기 지각 + SCR 의존",
    "▸ 최적 조합: 이른 분사시기 + 높은 분사압력 + 지각 점화시기 → Euro 7 대응",
]
multiline_box(s2, Inches(6.9), Inches(5.85), Inches(6.0), Inches(0.95),
              nox_summary, HL_G, GREEN, sz=8, co=DARK_GRAY)


# ============================================================
# 슬라이드 3: 핵심결과 (2) - 통합제어 + 2-stage TC
# ============================================================
s3 = prs.slides.add_slide(prs.slide_layouts[6])
s3.background.fill.solid(); s3.background.fill.fore_color.rgb = WHITE
title_bar(s3, "핵심결과 (2) - 통합제어 전략 및 고성능화",
          "운전영역별 연소-후처리 통합제어 · 2-Stage 터보차저 · 향후 로드맵")
footer(s3, 3)

# --- 좌측: 통합제어 ---
txt(s3, Inches(0.4), Inches(1.25), Inches(6), Inches(0.3),
    "5  운전영역별 연소-후처리 통합제어 전략", sz=14, b=True, co=DARK_BLUE)

ct = s3.shapes.add_table(5, 4, Inches(0.4), Inches(1.65), Inches(6.1), Inches(2.0)).table
ct.columns[0].width = Inches(1.3); ct.columns[1].width = Inches(1.7); ct.columns[2].width = Inches(1.5); ct.columns[3].width = Inches(1.6)
for j, h in enumerate(["운전영역", "열효율 지배인자", "NOx 제어", "제한 요소"]):
    cell(ct.cell(0, j), h, sz=9, b=True, co=WHITE, a=PP_ALIGN.CENTER, bg=TH_BG)
cd = [("중속 저부하", "공기과잉률 지배", "높은 λ", "-"),
      ("중속 중/고부하", "연소상+λ 동시", "높은 λ+지각", "NOx"),
      ("고속 저/중부하", "연소상+λ 동시", "낮은 λ 가능", "-"),
      ("고속 고부하", "연소상 지배", "낮은 λ+지각", "노킹/조기점화")]
for i, row in enumerate(cd):
    bg = ALT_BG if i % 2 == 0 else W_BG
    for j, v in enumerate(row):
        cell(ct.cell(i+1, j), v, sz=9, bg=bg, b=(j == 0), a=PP_ALIGN.CENTER)

# 최종 전략
txt(s3, Inches(0.4), Inches(3.8), Inches(6), Inches(0.25), "최종 전략 매트릭스", sz=12, b=True, co=DARK_BLUE)
st = s3.shapes.add_table(3, 3, Inches(0.4), Inches(4.1), Inches(6.1), Inches(1.0)).table
st.columns[0].width = Inches(1.0); st.columns[1].width = Inches(2.55); st.columns[2].width = Inches(2.55)
for j, h in enumerate(["부하", "중속", "고속"]):
    cell(st.cell(0, j), h, sz=10, b=True, co=WHITE, a=PP_ALIGN.CENTER, bg=TH_BG)
for i, (a, b, c) in enumerate([
    ("저부하", "높은 λ + 최적 연소상\n→ 최고 열효율", "낮은 λ + 최적 연소상"),
    ("고부하", "높은 λ + 지각 연소상\n→ NOx 억제", "낮은 λ + 지각 연소상\n→ 출력 확보"),
]):
    bg = HL_B if i == 0 else HL_Y
    for j, v in enumerate([a, b, c]):
        cell(st.cell(i+1, j), v, sz=9, bg=bg, b=(j == 0), a=PP_ALIGN.CENTER)

# 후처리
box(s3, Inches(0.4), Inches(5.3), Inches(2.85), Inches(0.5),
    "디젤: DOC→DPF→SCR\n(3단 구성)", MID_GRAY, sz=10)
txt(s3, Inches(3.3), Inches(5.35), Inches(0.4), Inches(0.4), "→", sz=20, b=True, co=GREEN, a=PP_ALIGN.CENTER)
box(s3, Inches(3.65), Inches(5.3), Inches(2.85), Inches(0.5),
    "수소: DOC→SCR\n(2단, DPF삭제+축소)", GREEN, sz=10)

box(s3, Inches(0.4), Inches(5.9), Inches(6.1), Inches(0.42),
    "전영역 TC후단 400℃↑ → SCR 99% → Euro 7 충족", HL_G, sz=10, tc=GREEN, b=True, lc=GREEN)

# --- 우측: 2-stage TC ---
txt(s3, Inches(6.9), Inches(1.25), Inches(5.5), Inches(0.3),
    "6  2-Stage 터보차저 적용 성과", sz=14, b=True, co=DARK_BLUE)

tt = s3.shapes.add_table(4, 4, Inches(6.9), Inches(1.65), Inches(6.0), Inches(1.5)).table
tt.columns[0].width = Inches(0.8); tt.columns[1].width = Inches(1.1); tt.columns[2].width = Inches(1.8); tt.columns[3].width = Inches(2.3)
for j, h in enumerate(["속도", "토크↑", "제한 요소", "최적화 전략"]):
    cell(tt.cell(0, j), h, sz=9, b=True, co=WHITE, a=PP_ALIGN.CENTER, bg=TH_BG)
td = [("저속", "+42.6%", "조기점화, 배기에너지↓", "점화시기 지각+적정 당량비"),
      ("중속", "+6.1%", "인터쿨러 전단온도, NOx", "고정 과급압 당량비 최적화"),
      ("고속", "+3.4%", "인터쿨러 전단온도, 조기점화", "NOx 마진 당량비 상향")]
for i, row in enumerate(td):
    is0 = (i == 0)
    bg = HL_Y if is0 else (ALT_BG if i % 2 == 0 else W_BG)
    cell(tt.cell(i+1, 0), row[0], sz=10, b=True, bg=bg, a=PP_ALIGN.CENTER)
    cell(tt.cell(i+1, 1), row[1], sz=13 if is0 else 11, b=True, bg=bg, co=ACCENT_ORANGE, a=PP_ALIGN.CENTER)
    cell(tt.cell(i+1, 2), row[2], sz=8, bg=bg, a=PP_ALIGN.CENTER)
    cell(tt.cell(i+1, 3), row[3], sz=8, bg=bg, a=PP_ALIGN.CENTER)

# 저속 메커니즘
txt(s3, Inches(6.9), Inches(3.3), Inches(5.5), Inches(0.25), "저속 토크 향상 메커니즘", sz=12, b=True, co=DARK_BLUE)
multiline_box(s3, Inches(6.9), Inches(3.6), Inches(6.0), Inches(1.4),
              ["▸ 점화시기 지각 → 배기에너지 확보 → 과급압 형성 (가장 효과적)",
               "   └ 한계: 실린더 압력/온도↑ → 조기점화 발생",
               "▸ 당량비 최적화:",
               "   └ 낮은 당량비: VGT duty↑ → 배기압↑ → 잔여가스↑ → 조기점화",
               "   └ 높은 당량비: 연소온도↑ → 조기점화 재발생",
               "   └ 적정 당량비: 조기점화 최소 + 열효율 최적"],
              LIGHT_GRAY, MEDIUM_BLUE, sz=9)

# 향후 로드맵
txt(s3, Inches(6.9), Inches(5.15), Inches(5.5), Inches(0.25), "향후 고성능화 로드맵", sz=12, b=True, co=DARK_BLUE)
rm = [("하이브리드화", "저속토크/반응성 보완", MEDIUM_BLUE),
      ("2단 인터쿨러", "중/고속 토크 추가향상", RGBColor(0x2E, 0x86, 0xAB)),
      ("Blowby 관리", "5%↑ H₂ 열점 관리", PURPLE)]
for i, (t, d, c) in enumerate(rm):
    x = Inches(6.9) + i * Inches(2.0)
    box(s3, x, Inches(5.45), Inches(1.9), Inches(0.35), t, c, sz=10)
    box(s3, x, Inches(5.85), Inches(1.9), Inches(0.42), d, LIGHT_GRAY, sz=9, tc=DARK_GRAY, b=False, lc=c)

# 최하단 키메시지
box(s3, Inches(0.4), Inches(6.5), Inches(12.5), Inches(0.42),
    "300kW 출력 + 41.4% 열효율 + Euro 7 충족 + 저속토크 42.6%↑ → 친환경 대형 동력원 상용화 기반 완성",
    RGBColor(0xFD, 0xF2, 0xE9), sz=11, tc=ACCENT_ORANGE, b=True, lc=ACCENT_ORANGE)

out = r"C:\Users\hwpar\OneDrive\퇴직준비세미나\Projects\요약자료\output\300kW_수소엔진_발표자료_v3.pptx"
prs.save(out)
print(f"PPTX v3 생성 완료: {out}")
