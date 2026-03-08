"""300kW 수소엔진 최종 발표자료 PPTX 생성 스크립트"""
import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 와이드스크린 16:9
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 색상 정의
DARK_BLUE = RGBColor(0x00, 0x2B, 0x5C)
MEDIUM_BLUE = RGBColor(0x00, 0x4E, 0x92)
LIGHT_BLUE = RGBColor(0x00, 0x7B, 0xC0)
ACCENT_ORANGE = RGBColor(0xE8, 0x6C, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MID_GRAY = RGBColor(0x66, 0x66, 0x66)
TABLE_HEADER_BG = RGBColor(0x00, 0x2B, 0x5C)
TABLE_ALT_BG = RGBColor(0xE8, 0xF0, 0xF8)
TABLE_WHITE_BG = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x00, 0x7A, 0x33)
RED_ACCENT = RGBColor(0xC0, 0x39, 0x2B)


def add_slide_background(slide, color=DARK_BLUE):
    """슬라이드 배경색 설정"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_bar(slide, title_text, subtitle_text=None):
    """상단 타이틀 바"""
    # 타이틀 바 배경
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(1.1)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    # 타이틀 텍스트
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12), Inches(0.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    if subtitle_text:
        txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.7), Inches(12), Inches(0.35))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle_text
        p2.font.size = Pt(14)
        p2.font.color.rgb = RGBColor(0xAA, 0xCC, 0xEE)

    # 하단 악센트 라인
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(1.1),
        Inches(13.333), Inches(0.04)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_ORANGE
    line.line.fill.background()


def add_footer(slide, page_num, total=3):
    """하단 페이지 번호"""
    txBox = slide.shapes.add_textbox(Inches(11.5), Inches(7.1), Inches(1.5), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{page_num} / {total}"
    p.font.size = Pt(10)
    p.font.color.rgb = MID_GRAY
    p.alignment = PP_ALIGN.RIGHT

    # 기관명
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(4), Inches(0.3))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "한국기계연구원(KIMM) | 산업기술혁신사업"
    p2.font.size = Pt(10)
    p2.font.color.rgb = MID_GRAY


def set_cell_format(cell, text, font_size=10, bold=False, color=DARK_GRAY,
                    alignment=PP_ALIGN.LEFT, bg_color=None):
    """테이블 셀 포맷 설정"""
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if bg_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg_color


def add_rounded_box(slide, left, top, width, height, text, bg_color, font_size=11,
                    text_color=WHITE, bold=True):
    """둥근 모서리 박스"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(6)
    tf.margin_top = Pt(3)
    tf.margin_bottom = Pt(3)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.CENTER
    return shape


# ============================================================
# 슬라이드 1: 1~4차년도 요약장표
# ============================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_slide_background(slide1, WHITE)
add_title_bar(slide1, "300kW급 수소 엔진 개발 - 연차별 추진 요약",
              "공동연구개발기관: 한국기계연구원 | 1단계 1~4차년도")
add_footer(slide1, 1)

# 4개 연차 카드 배치
card_data = [
    {
        "year": "1차년도",
        "theme": "시스템 구축 · 기초 성능",
        "items": [
            "동력계 엔진 성능 평가 시스템 구축\n(460kW AC동력계)",
            "수소 연료 공급 설비 증축\n(35kg 저장, 60bar 대응)",
            "베이스 천연가스 엔진(GX12)\n성능 DB 구축",
            "수소 엔진 기초 시험\n→ 천연가스 대비 열효율 2.8%p↑",
        ],
        "highlight": "열효율 2.8%p 향상",
        "color": RGBColor(0x00, 0x6B, 0x9F),
    },
    {
        "year": "2차년도",
        "theme": "하드웨어 선정 · 요소기술",
        "items": [
            "냉형 점화플러그 적용\n→ 이른 조기점화 억제",
            "포트분사 인젝터 사양 선정\n(분사 인가기간 <10ms)",
            "최적 압축비 11:1 선정\n→ 열효율 ~1%p 향상",
            "이상연소 제어 전략 수립\n(공기과잉률 + EGR)",
        ],
        "highlight": "압축비 11:1 최적화",
        "color": RGBColor(0x00, 0x4E, 0x92),
    },
    {
        "year": "3차년도",
        "theme": "NOx 저감 · 후처리시스템",
        "items": [
            "NOx 저감 연소 제어 전략\n(분사/점화시기, 공기과잉률)",
            "저부하 EGR 전략\n→ 펌핑손실 저감",
            "후처리시스템 간소화\nDOC + SCR (DPF 삭제)",
            "엔진-후처리 연계 시험\n→ Euro 7 배출규제 충족",
        ],
        "highlight": "Euro 7 충족 달성",
        "color": RGBColor(0x2E, 0x86, 0xAB),
    },
    {
        "year": "4차년도",
        "theme": "통합제어 · 고성능화",
        "items": [
            "운전영역별 연소-후처리\n통합제어 전략 도출",
            "주행 사이클 대응\n(WHSC/WHTC/NRSC/NRTC)",
            "2-stage 터보차저 적용\n→ 저속토크 42.6%↑",
            "고성능 수소엔진\n개발 로드맵 제시",
        ],
        "highlight": "저속토크 42.6% 향상",
        "color": RGBColor(0x00, 0x2B, 0x5C),
    },
]

card_width = Inches(3.0)
card_gap = Inches(0.15)
start_x = Inches(0.45)
card_top = Inches(1.4)

for i, card in enumerate(card_data):
    x = start_x + i * (card_width + card_gap)

    # 연차 헤더
    add_rounded_box(slide1, x, card_top, card_width, Inches(0.45),
                    card["year"], card["color"], font_size=14, bold=True)

    # 테마
    add_rounded_box(slide1, x, card_top + Inches(0.5), card_width, Inches(0.35),
                    card["theme"], RGBColor(0xEE, 0xEE, 0xEE), font_size=11,
                    text_color=DARK_GRAY, bold=True)

    # 항목들
    item_top = card_top + Inches(1.0)
    for j, item in enumerate(card["items"]):
        box = slide1.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x + Inches(0.05), item_top + j * Inches(1.15),
            card_width - Inches(0.1), Inches(1.05)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_GRAY
        box.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
        box.line.width = Pt(0.5)

        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(8)
        tf.margin_right = Pt(8)
        tf.margin_top = Pt(5)
        tf.margin_bottom = Pt(5)
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(10)
        p.font.color.rgb = DARK_GRAY
        p.line_spacing = Pt(14)

    # 핵심 성과 하이라이트
    highlight_top = item_top + 4 * Inches(1.15) + Inches(0.1)
    add_rounded_box(slide1, x, highlight_top, card_width, Inches(0.4),
                    f"★ {card['highlight']}", ACCENT_ORANGE, font_size=11, bold=True)


# ============================================================
# 슬라이드 2: 핵심결과 (1) - 엔진 성능
# ============================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide2, WHITE)
add_title_bar(slide2, "핵심결과 (1) - 엔진 성능 달성",
              "출력 · 열효율 · 2-Stage 터보차저 · 이상연소 제어")
add_footer(slide2, 2)

# 좌측: 출력 및 열효율 테이블
section_label1 = slide2.shapes.add_textbox(Inches(0.5), Inches(1.35), Inches(5), Inches(0.35))
tf = section_label1.text_frame
p = tf.paragraphs[0]
p.text = "출력 및 열효율 달성"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = DARK_BLUE

# 성능 테이블
perf_table = slide2.shapes.add_table(5, 2, Inches(0.5), Inches(1.8), Inches(5.8), Inches(2.6)).table
perf_table.columns[0].width = Inches(2.5)
perf_table.columns[1].width = Inches(3.3)

headers = ["항목", "달성 결과"]
data = [
    ["최고출력", "300kW (2,000 rpm)"],
    ["최고 열효율", "41.4% (1,200rpm, 100% 부하, CR 11:1)"],
    ["연소효율", "전 영역 99% 이상"],
    ["천연가스 대비 개선", "+2.8%p (연료 변경만으로 달성)"],
]

for j, h in enumerate(headers):
    set_cell_format(perf_table.cell(0, j), h, font_size=11, bold=True,
                    color=WHITE, alignment=PP_ALIGN.CENTER, bg_color=TABLE_HEADER_BG)

for i, row in enumerate(data):
    bg = TABLE_ALT_BG if i % 2 == 0 else TABLE_WHITE_BG
    set_cell_format(perf_table.cell(i + 1, 0), row[0], font_size=11, bold=True,
                    bg_color=bg)
    set_cell_format(perf_table.cell(i + 1, 1), row[1], font_size=11,
                    bg_color=bg, color=DARK_BLUE, bold=True)

# 우측: 2-Stage 터보차저
section_label2 = slide2.shapes.add_textbox(Inches(7), Inches(1.35), Inches(5.5), Inches(0.35))
tf2 = section_label2.text_frame
p2 = tf2.paragraphs[0]
p2.text = "2-Stage 터보차저 적용 성과"
p2.font.size = Pt(16)
p2.font.bold = True
p2.font.color.rgb = DARK_BLUE

tc_table = slide2.shapes.add_table(4, 3, Inches(7), Inches(1.8), Inches(5.8), Inches(2.6)).table
tc_table.columns[0].width = Inches(1.3)
tc_table.columns[1].width = Inches(1.8)
tc_table.columns[2].width = Inches(2.7)

tc_headers = ["엔진 속도", "토크 향상", "제한 요소"]
tc_data = [
    ["저속", "+42.6%", "낮은 배기에너지, 조기점화"],
    ["중속", "+6.1%", "인터쿨러 전단 온도, NOx"],
    ["고속", "+3.4%", "인터쿨러 전단 온도, 조기점화"],
]

for j, h in enumerate(tc_headers):
    set_cell_format(tc_table.cell(0, j), h, font_size=11, bold=True,
                    color=WHITE, alignment=PP_ALIGN.CENTER, bg_color=TABLE_HEADER_BG)

for i, row in enumerate(tc_data):
    bg = TABLE_ALT_BG if i % 2 == 0 else TABLE_WHITE_BG
    set_cell_format(tc_table.cell(i + 1, 0), row[0], font_size=11,
                    alignment=PP_ALIGN.CENTER, bg_color=bg, bold=True)
    set_cell_format(tc_table.cell(i + 1, 1), row[1], font_size=14,
                    alignment=PP_ALIGN.CENTER, bg_color=bg, bold=True,
                    color=ACCENT_ORANGE)
    set_cell_format(tc_table.cell(i + 1, 2), row[2], font_size=10,
                    bg_color=bg)

# 하단: 이상연소 억제 전략
section_label3 = slide2.shapes.add_textbox(Inches(0.5), Inches(4.7), Inches(12), Inches(0.35))
tf3 = section_label3.text_frame
p3 = tf3.paragraphs[0]
p3.text = "이상연소 억제 전략 체계"
p3.font.size = Pt(16)
p3.font.bold = True
p3.font.color.rgb = DARK_BLUE

abnormal_data = [
    ("역화 억제", "공기과잉률 증가 (λ ≥ 2.0)", RGBColor(0xC0, 0x39, 0x2B)),
    ("조기점화 억제", "냉형 점화플러그 + 분사시기/압력 최적화", RGBColor(0xE6, 0x7E, 0x22)),
    ("노킹 억제", "최적 압축비(11:1) + 점화시기 지각", RGBColor(0x27, 0xAE, 0x60)),
    ("종합 제어", "EGR 적용(저부하) + 공기과잉률 조절(전 영역)", MEDIUM_BLUE),
]

for i, (label, desc, color) in enumerate(abnormal_data):
    x_pos = Inches(0.5) + i * Inches(3.15)
    # 라벨
    add_rounded_box(slide2, x_pos, Inches(5.15), Inches(1.8), Inches(0.4),
                    label, color, font_size=12, bold=True)
    # 화살표 영역
    arrow_box = slide2.shapes.add_textbox(x_pos + Inches(1.9), Inches(5.18), Inches(0.3), Inches(0.35))
    tf_a = arrow_box.text_frame
    p_a = tf_a.paragraphs[0]
    p_a.text = "→"
    p_a.font.size = Pt(16)
    p_a.font.color.rgb = color
    p_a.font.bold = True
    # 설명
    desc_box = slide2.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        x_pos, Inches(5.65), Inches(3.0), Inches(0.7)
    )
    desc_box.fill.solid()
    desc_box.fill.fore_color.rgb = LIGHT_GRAY
    desc_box.line.color.rgb = color
    desc_box.line.width = Pt(1.5)
    tf_d = desc_box.text_frame
    tf_d.word_wrap = True
    tf_d.margin_left = Pt(8)
    tf_d.margin_right = Pt(8)
    p_d = tf_d.paragraphs[0]
    p_d.text = desc
    p_d.font.size = Pt(10)
    p_d.font.color.rgb = DARK_GRAY
    p_d.alignment = PP_ALIGN.CENTER

# 하단 키 메시지
key_msg = slide2.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.45)
)
key_msg.fill.solid()
key_msg.fill.fore_color.rgb = RGBColor(0xFD, 0xF2, 0xE9)
key_msg.line.color.rgb = ACCENT_ORANGE
key_msg.line.width = Pt(1)
tf_k = key_msg.text_frame
tf_k.word_wrap = True
p_k = tf_k.paragraphs[0]
p_k.text = ("핵심: 300kW 출력 달성 + 열효율 41.4% + "
            "이상연소 체계적 억제 → 상용화 기반 확보")
p_k.font.size = Pt(12)
p_k.font.bold = True
p_k.font.color.rgb = ACCENT_ORANGE
p_k.alignment = PP_ALIGN.CENTER


# ============================================================
# 슬라이드 3: 핵심결과 (2) - 배출 및 통합 전략
# ============================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide3, WHITE)
add_title_bar(slide3, "핵심결과 (2) - 배출 특성 및 통합 제어 전략",
              "배출가스 · 후처리시스템 · 운전영역별 전략 · 향후 과제")
add_footer(slide3, 3)

# 좌측 상단: 배출가스 특성
section_label4 = slide3.shapes.add_textbox(Inches(0.5), Inches(1.35), Inches(6), Inches(0.35))
tf4 = section_label4.text_frame
p4 = tf4.paragraphs[0]
p4.text = "배출가스 특성 및 대응 전략"
p4.font.size = Pt(16)
p4.font.bold = True
p4.font.color.rgb = DARK_BLUE

emit_table = slide3.shapes.add_table(5, 3, Inches(0.5), Inches(1.8), Inches(6.0), Inches(2.5)).table
emit_table.columns[0].width = Inches(1.0)
emit_table.columns[1].width = Inches(2.2)
emit_table.columns[2].width = Inches(2.8)

emit_headers = ["배출물", "수소 엔진 특성", "대응 전략"]
emit_data = [
    ["NOx", "유일한 유해 배출물", "공기과잉률↑ + Urea SCR (99%)"],
    ["CO/CO₂", "제로 수준", "무탄소 연료 효과"],
    ["THC", "오일연소 일부 발생", "오일유량/피스톤 간극 최적화"],
    ["PM", "제로 수준", "DPF 불필요 → 후처리 간소화"],
]

for j, h in enumerate(emit_headers):
    set_cell_format(emit_table.cell(0, j), h, font_size=11, bold=True,
                    color=WHITE, alignment=PP_ALIGN.CENTER, bg_color=TABLE_HEADER_BG)

for i, row in enumerate(emit_data):
    bg = TABLE_ALT_BG if i % 2 == 0 else TABLE_WHITE_BG
    set_cell_format(emit_table.cell(i + 1, 0), row[0], font_size=11,
                    alignment=PP_ALIGN.CENTER, bg_color=bg, bold=True)
    set_cell_format(emit_table.cell(i + 1, 1), row[1], font_size=10, bg_color=bg)
    set_cell_format(emit_table.cell(i + 1, 2), row[2], font_size=10, bg_color=bg,
                    color=MEDIUM_BLUE, bold=True)

# 우측 상단: 운전영역별 통합 전략
section_label5 = slide3.shapes.add_textbox(Inches(7), Inches(1.35), Inches(5.5), Inches(0.35))
tf5 = section_label5.text_frame
p5 = tf5.paragraphs[0]
p5.text = "운전영역별 통합 제어 전략"
p5.font.size = Pt(16)
p5.font.bold = True
p5.font.color.rgb = DARK_BLUE

ctrl_table = slide3.shapes.add_table(3, 3, Inches(7), Inches(1.8), Inches(5.8), Inches(1.5)).table
ctrl_table.columns[0].width = Inches(1.2)
ctrl_table.columns[1].width = Inches(2.3)
ctrl_table.columns[2].width = Inches(2.3)

ctrl_headers = ["엔진 부하", "중속 전략", "고속 전략"]
ctrl_data = [
    ["저부하", "높은 λ + 최적 연소상", "낮은 λ + 최적 연소상"],
    ["고부하", "높은 λ + 지각 연소상", "낮은 λ + 지각 연소상"],
]

for j, h in enumerate(ctrl_headers):
    set_cell_format(ctrl_table.cell(0, j), h, font_size=11, bold=True,
                    color=WHITE, alignment=PP_ALIGN.CENTER, bg_color=TABLE_HEADER_BG)

for i, row in enumerate(ctrl_data):
    bg = TABLE_ALT_BG if i % 2 == 0 else TABLE_WHITE_BG
    set_cell_format(ctrl_table.cell(i + 1, 0), row[0], font_size=11,
                    alignment=PP_ALIGN.CENTER, bg_color=bg, bold=True)
    set_cell_format(ctrl_table.cell(i + 1, 1), row[1], font_size=10,
                    alignment=PP_ALIGN.CENTER, bg_color=bg)
    set_cell_format(ctrl_table.cell(i + 1, 2), row[2], font_size=10,
                    alignment=PP_ALIGN.CENTER, bg_color=bg)

# 전략 설명 박스
strategy_box = slide3.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(7), Inches(3.4), Inches(5.8), Inches(0.8)
)
strategy_box.fill.solid()
strategy_box.fill.fore_color.rgb = RGBColor(0xE8, 0xF0, 0xF8)
strategy_box.line.color.rgb = MEDIUM_BLUE
strategy_box.line.width = Pt(1)
tf_s = strategy_box.text_frame
tf_s.word_wrap = True
tf_s.margin_left = Pt(10)
tf_s.margin_right = Pt(10)
tf_s.margin_top = Pt(6)
p_s = tf_s.paragraphs[0]
p_s.text = ("▸ 중속: 높은 Engine-out NOx → 높은 공기과잉률 적용\n"
            "▸ 고속: 낮은 Engine-out NOx → 낮은 공기과잉률로 출력 확보")
p_s.font.size = Pt(10)
p_s.font.color.rgb = DARK_BLUE
p_s.line_spacing = Pt(16)

# 하단 좌측: 후처리시스템 간소화
section_label6 = slide3.shapes.add_textbox(Inches(0.5), Inches(4.6), Inches(6), Inches(0.35))
tf6 = section_label6.text_frame
p6 = tf6.paragraphs[0]
p6.text = "후처리시스템 간소화"
p6.font.size = Pt(16)
p6.font.bold = True
p6.font.color.rgb = DARK_BLUE

# 디젤 vs 수소 비교
diesel_box = add_rounded_box(slide3, Inches(0.5), Inches(5.1), Inches(5.8), Inches(0.5),
                             "기존 디젤: DOC → DPF → SCR (3단 구성)", MID_GRAY, font_size=12)

arrow_down = slide3.shapes.add_textbox(Inches(3.0), Inches(5.6), Inches(0.8), Inches(0.3))
tf_ad = arrow_down.text_frame
p_ad = tf_ad.paragraphs[0]
p_ad.text = "▼"
p_ad.font.size = Pt(16)
p_ad.font.color.rgb = GREEN
p_ad.font.bold = True
p_ad.alignment = PP_ALIGN.CENTER

h2_box = add_rounded_box(slide3, Inches(0.5), Inches(5.9), Inches(5.8), Inches(0.5),
                         "수소 엔진: DOC → SCR (2단 구성, DPF 삭제 + 사이즈 축소)",
                         GREEN, font_size=12)

# 하단 우측: 향후 개발 방향
section_label7 = slide3.shapes.add_textbox(Inches(7), Inches(4.6), Inches(5.5), Inches(0.35))
tf7 = section_label7.text_frame
p7 = tf7.paragraphs[0]
p7.text = "향후 개발 방향"
p7.font.size = Pt(16)
p7.font.bold = True
p7.font.color.rgb = DARK_BLUE

future_items = [
    ("하이브리드화", "디젤 대비 낮은 저속토크\n및 반응성 보완"),
    ("2단 인터쿨러", "Mid charge air cooling\n→ 중/고속 추가 토크 향상"),
    ("Blowby 관리", "5% 이상 수소 농도\n→ 열점 관리 필수"),
]

for i, (title, desc) in enumerate(future_items):
    x_pos = Inches(7) + i * Inches(1.95)
    add_rounded_box(slide3, x_pos, Inches(5.1), Inches(1.85), Inches(0.4),
                    title, MEDIUM_BLUE, font_size=11, bold=True)
    desc_box = slide3.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        x_pos, Inches(5.6), Inches(1.85), Inches(0.85)
    )
    desc_box.fill.solid()
    desc_box.fill.fore_color.rgb = LIGHT_GRAY
    desc_box.line.color.rgb = MEDIUM_BLUE
    desc_box.line.width = Pt(0.5)
    tf_d = desc_box.text_frame
    tf_d.word_wrap = True
    tf_d.margin_left = Pt(6)
    tf_d.margin_right = Pt(6)
    tf_d.margin_top = Pt(4)
    p_d = tf_d.paragraphs[0]
    p_d.text = desc
    p_d.font.size = Pt(9)
    p_d.font.color.rgb = DARK_GRAY
    p_d.line_spacing = Pt(13)
    p_d.alignment = PP_ALIGN.CENTER

# 최하단 키 메시지
key_msg2 = slide3.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.45)
)
key_msg2.fill.solid()
key_msg2.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
key_msg2.line.color.rgb = GREEN
key_msg2.line.width = Pt(1)
tf_k2 = key_msg2.text_frame
tf_k2.word_wrap = True
p_k2 = tf_k2.paragraphs[0]
p_k2.text = ("핵심: 무탄소 배출 + Euro 7 충족 + 후처리 간소화(DOC+SCR) "
             "→ 친환경 대형 동력원 실용화 기반 완성")
p_k2.font.size = Pt(12)
p_k2.font.bold = True
p_k2.font.color.rgb = GREEN
p_k2.alignment = PP_ALIGN.CENTER


# 저장
output_path = r"C:\Users\hwpar\OneDrive\퇴직준비세미나\Projects\요약자료\output\300kW_수소엔진_발표자료.pptx"
prs.save(output_path)
import sys
sys.stdout.reconfigure(encoding='utf-8')
print(f"PPTX 생성 완료: {output_path}")
